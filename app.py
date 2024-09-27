import string 
import csv
from io import StringIO
from pptx import Presentation
import docx2txt
import PyPDF2
import spacy
import pandas as pd 
import numpy as np
import nltk 
import re
import openpyxl
from nltk.stem import WordNetLemmatizer
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
from gensim.parsing.preprocessing import STOPWORDS as SW
nltk.download('stopwords')
nltk.download('wordnet')
nltk.download('omw-1.4')
nltk.download('averaged_perceptron_tagger')
from nltk.corpus import wordnet
import networkx as nx
from networkx.algorithms.shortest_paths import weighted
import glob

class pdfReader:
    
    def __init__(self, file_path: str) -> str:
        self.file_path = file_path
    
    def PDF_one_pager(self) -> str:
        """A function returns a one line string of the 
            pdf.
            
            Returns:
            one_page_pdf (str): A one line string of the pdf.
        
        """
        content = ""
        p = open(self.file_path, "rb")
        pdf = PyPDF2.PdfReader(p)
        num_pages = len(pdf.pages)
        for i in range(0, num_pages):
            content += pdf.pages[i].extract_text() + "\n"
        content = " ".join(content.replace(u"\xa0", " ").strip().split())
        page_number_removal = r"\d{1,3} of \d{1,3}"
        page_number_removal_pattern = re.compile(page_number_removal, re.IGNORECASE)
        content = re.sub(page_number_removal_pattern, '',content)
        
        return content
  
    def pdf_reader(self) -> str:
        """A function that can read .pdf formatted files 
            and returns a python readable pdf.
            
            Returns:
            read_pdf: A python readable .pdf file.
        """
        opener = open(self.file_path,'rb')
        read_pdf = PyPDF2.PdfFileReader(opener)
    
        return read_pdf
  
  
    def pdf_info(self) -> dict:
        """A function which returns an information dictionary of a 
        pdf.
        
        Returns:
        dict(pdf_info_dict): A dictionary containing the meta
        data of the object.
        """
        opener = open(self.file_path,'rb')
        read_pdf = PyPDF2.PdfFileReader(opener)
        pdf_info_dict = {}
        for key,value in read_pdf.documentInfo.items():
            pdf_info_dict[re.sub('/',"",key)] = value
        return pdf_info_dict
  
    def pdf_dictionary(self) -> dict:
        """A function which returns a dictionary of 
            the object where the keys are the pages
            and the text within the pages are the 
            values.

            Returns:
            dict(pdf_dict): A dictionary pages and text.
        """
        opener = open(self.file_path,'rb')
  
        read_pdf = PyPDF2.PdfReader(opener)
        length = read_pdf.pages
        pdf_dict = {}
        for i in range(length):
            page = read_pdf.getPage(i)
            text = page.extract_text()
            pdf_dict[i] = text
            return pdf_dict
          
class xlsxReader:
    
    def __init__(self, file_path: str) -> str:
        self.file_path = file_path
        
    def xlsx_text(self):
      """A function which returns a string of an 
         excel document.
         
          Returns:
          text(str): String of text of a document.
      """
      inputExcelFile = self.file_path
      text = str()
      wb = openpyxl.load_workbook(inputExcelFile)
      for sn in wb.sheetnames:
        excelFile = pd.read_excel(inputExcelFile, engine = 'openpyxl', sheet_name = sn)
        excelFile.to_csv("ResultCsvFile.csv", index = None, header=True)

        with open("ResultCsvFile.csv", "r") as csvFile: 
          lines = csvFile.read().split(",") # "\r\n" if needed
          for val in lines:
            if val != '':
              text += val + ' '
          text = text.replace('\ufeff', '')
          text = text.replace('\n', ' ')
      return text
    
class csvReader:
    
    def __init__(self, file_path: str) -> str:
        self.file_path = file_path
        
    def csv_text(self):
      """A function which returns a string of an 
         csv document.
         
          Returns:
          text(str): String of text of a document.
      """
      text = str()
      with open(self.file_path, "r") as csvFile: 
        lines = csvFile.read().split(",") # "\r\n" if needed
        for val in lines:
          text += val + ' '
        text = text.replace('\ufeff', '')
        text = text.replace('\n', ' ')
      return text

class pptReader:
    
    def __init__(self, file_path: str) -> str:
        self.file_path = file_path
        
    def ppt_text(self):
      """A function which returns a string of an 
        Mirocsoft PowerPoint document.
        
        Returns:
        text(str): String of text of a document.
    """
      prs = Presentation(self.file_path)
      text = str()
      for slide in prs.slides:
        for shape in slide.shapes:
          if not shape.has_text_frame:
              continue
          for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
              text += ' ' + run.text
                  
      return text
    
class wordDocReader:
  def __init__(self, file_path: str) -> str:
    self.file_path = file_path
    
  def word_reader(self):
    """A function which returns a string of an 
          Microsoft Word document.
          
          Returns:
          text(str): String of text of a document.
      """
    text = docx2txt.process(self.file_path)
    text = text.replace('\n', ' ')
    text = text.replace('\xa0', ' ')
    text = text.replace('\t', ' ')
    return text 
  
class dataprocessor:
  def __init__(self):
    return
  
  @staticmethod
  def get_wordnet_pos(text: str) -> str:
    """Map POS tag to first character lemmatize() accepts
    Inputs:
    text(str): A string of text

    Returns:
    tag_dict(dict): A dictionary of tags
    """
    tag = nltk.pos_tag([text])[0][1][0].upper()
    tag_dict = {"J": wordnet.ADJ,
                "N": wordnet.NOUN,
                "V": wordnet.VERB,
                "R": wordnet.ADV}

    return tag_dict.get(tag, wordnet.NOUN)
  
  @staticmethod
  def preprocess(text: str):
    """A function that prepoccesses text through the
    steps of Natural Language Processing (NLP).
      Inputs:
      text(str): A string of text

      Returns:
      text(str): A processed string of text
      """
    #lowercase
    text = text.lower()
    
    #punctuation removal
    text = "".join([i for i in text if i not in string.punctuation])
    
    #Digit removal (Only for ALL numeric numbers)
    text = [x for x in text.split(' ') if x.isnumeric() == False]
    
    #Stop removal
    stopwords = nltk.corpus.stopwords.words('english')
    custom_stopwords = ['\n','\n\n', '&', ' ', '.', '-', '$', '@']
    stopwords.extend(custom_stopwords)

    text = [i for i in text if i not in stopwords]
    text = ' '.join(word for word in text)
    
    #lemmanization
    lm = WordNetLemmatizer()
    text = [lm.lemmatize(word, dataprocessor.get_wordnet_pos(word)) for word in text.split(' ')]
    text = ' '.join(word for word in text)
    
    text = re.sub(' +', ' ',text)
    
    return text


  @staticmethod
  def data_reader(list_file_names,file_dict=dict()):
    """A function that reads in the data from a directory of files.

    Inputs:
    list_file_names(list): List of the filepaths in a directory.

    Returns:
    text_list (list): A list where each value is a string of the text
    for each file in the directory
    file_dict(dict): Dictionary where the keys are the filename and the values
    are the information found within each given file
    """

    text_list = []
    reader = dataprocessor()
    for file in list_file_names:
      temp = file.split('.')
      filetype = temp[-1]
      if filetype == "pdf":
        file_pdf = pdfReader(file)
        text = file_pdf.PDF_one_pager()

      elif filetype == "docx":
        word_doc_reader = wordDocReader(file)
        text = word_doc_reader.word_reader()

      elif filetype == "pptx" or filetype == 'ppt':
        ppt_reader = pptReader(file)
        text = ppt_reader.ppt_text()
        
      elif filetype == "csv":
        csv_reader = csvReader(file)
        text = csv_reader.csv_text()
      
      elif filetype == 'xlsx':
        xl_reader = xlsxReader(file)
        text = xl_reader.xlsx_text()
      else:
        print('File type {} not supported!'.format(filetype))
        continue
      
      text = reader.preprocess(text)
      text_list.append(text)
      for i,file in enumerate(list_file_names):
        file_dict[i] = (file, file.split('/')[-1])
    return text_list, file_dict

  @staticmethod
  def database_processor(file_dict,text_list: list):
    """A function that transforms the text of each file within the 
    database into a vector.

    Inputs:
    file_dixt(dict): Dictionary where the keys are the filename and the values
    are the information found within each given file
    text_list (list): A list where each value is a string of the text
    for each file in the directory

    Returns:
    list_dense(list): A list of the files' text turned into vectors.
    vectorizer: The vectorizor used to transform the strings of text
    file_vector_dict(dict): A dictionary where the file names are the keys
    and the vectors of each files' text are the values.
    """
    file_vector_dict = dict()
    vectorizer = TfidfVectorizer()
    vectors = vectorizer.fit_transform(text_list)
    feature_names = vectorizer.get_feature_names_out()
    matrix = vectors.todense()
    list_dense = matrix.tolist()
    for i in range(len(list_dense)):
      file_vector_dict[file_dict[i][1]] = list_dense[i]

    return list_dense, vectorizer, file_vector_dict

  @staticmethod
  def input_processor(text, TDIF_vectorizor):
    """A function accepts a string of text and vectorizes the text using a 
     TDIF vectorizoer.

    Inputs:
    text(str): A string of text
    TDIF_vectorizor: A pretrained vectorizor

    Returns:
    words(list): A list of the input text in vectored form.
    """
    words = ''
    total_words = len(text.split(' '))
    for word in text.split(' '):
      words += (word + ' ') * total_words
      total_words -= 1

    words = [words[:-1]]
    words = TDIF_vectorizor.transform(words)
    words = words.todense()
    words = words.tolist()
    return words

  @staticmethod
  def similarity_checker(vector_1, vector_2):
    """A function accepts two vectors and computes their cosine similarity.

    Inputs:
    vector_1(int): A numerical vector
    vector_2(int): A numerical vector

    Returns:
    cosine_similarity([vector_1], vector_2) (int): Cosine similarity score
    """
    vectors = [vector_1, vector_2]
    for vec in vectors:
      if np.ndim(vec) == 1:
        vec = np.expand_dims(vec, axis=0)
    return cosine_similarity([vector_1], vector_2)

  @staticmethod
  def recommender(vector_file_list,query_vector, file_dict):
    """A function accepts a list of vectors, query vectors, and a dictionary
    pertaining to the list of vectors with their original values and file names.

    Inputs:
    vector_file_list(list): A list of vectors
    query_vector(int): A numerical vector
    file_dict(dict): A dictionary of filenames and text relating to the list
    of vectors

    Returns:
    final_recommendation (list): A list of the final recommended files
    similarity_list[:len(final_recommendation)] (list): A list of the similarity
    scores of the final recommendations.
    """
    similarity_list = []
    score_dict = dict()
    for i,file_vector in enumerate(vector_file_list):
      x = dataprocessor.similarity_checker(file_vector, query_vector)
      score_dict[file_dict[i][1]] = (x[0][0])
      similarity_list.append(x)
    similarity_list = sorted(similarity_list, reverse = True)
    #Recommends the top 20%
    recommended = sorted(score_dict.items(), 
                  key=lambda x:-x[1])[:int(np.round(.5*len(similarity_list)))]
  
    final_recommendation = []
    for i in range(len(recommended)):
      final_recommendation.append(recommended[i][0])
    #add in graph for greater than 3 recommendationa
    return final_recommendation, similarity_list[:len(final_recommendation)]

  @staticmethod
  def ranker(recommendation_val, file_vec_dict):
    """A function accepts a list of recommendaton values and a dictionary
    files wihin the databse and their vectors.

    Inputs:
    reccomendation_val(list): A list of recommendations found through cosine
    similarity
    file_vec_dic(dict): A dictionary of the filenames as keys and their
    text in vectors as the values.

    Returns:
    ec_recommended(list): A list of the top 20% recommendations found using the 
    eigenvector centrality algorithm.
    """
    my_graph = nx.Graph()
    for i in range(len(recommendation_val)):
      file_1 = recommendation_val[i]
      for j in range(len(recommendation_val)):
        file_2 = recommendation_val[j]

        if i != j:
          #Calculate sim_score between two values (weight)
          edge_dist = cosine_similarity([file_vec_dict[recommendation_val[i]]],[file_vec_dict[recommendation_val[j]]])
          #add an edge from file 1 to file 2 with the weight 
          my_graph.add_edge(file_1, file_2, weight=edge_dist)

    #Pagerank the graph  ]    
    rec = nx.eigenvector_centrality(my_graph)
    #Takes 20% of the values
    ec_recommended = sorted(rec.items(), key=lambda x:-x[1])[:int(np.round(len(rec)))]

    return ec_recommended

  @staticmethod
  def weighted_final_rank(sim_list,ec_recommended,final_recommendation):
    """A function accepts a list of similiarity values found through 
      cosine similairty, similarities found through eigenvector centrality,
      and the final recommendations produced by cosine similarity.

        Inputs:
        sim_list(list): A list of all of the similarity values for the files
        within the database.
        ec_recommended(list): A list of the top 20% recommendations found using the 
        eigenvector centrality algorithm.
        final_recommendation (list): A list of the final recommendations found
        by using cosine similarity.

        Returns:
        weighted_final_recommend(list): A list of the final recommendations for 
        the files in the database.
        """
    final_dict = dict()

    for i in range(len(sim_list)):
      val = (.8*sim_list[final_recommendation.index(ec_recommendation[i][0])].squeeze()) + (.2 * ec_recommendation[i][1])
      final_dict[ec_recommendation[i][0]] = val

    weighted_final_recommend = sorted(final_dict.items(), key=lambda x:-x[1])[:int(np.round(len(final_dict)))]

    return weighted_final_recommend
  
 path = '/content/drive/MyDrive/database/'
db = [f for f in glob.glob(path + '*')]

research_documents, file_dictionary = dataprocessor.data_reader(db)
list_files, vectorizer, file_vec_dict = dataprocessor.database_processor(file_dictionary,research_documents)
query = 'Machine Learning'
query = dataprocessor.preprocess(query)
query = dataprocessor.input_processor(query, vectorizer)
recommendation, sim_list = dataprocessor.recommender(list_files,query, file_dictionary)
ec_recommendation = dataprocessor.ranker(recommendation, file_vec_dict)
final_weighted_recommended = dataprocessor.weighted_final_rank(sim_list,ec_recommendation,  recommendation)
print(final_weighted_recommended)
